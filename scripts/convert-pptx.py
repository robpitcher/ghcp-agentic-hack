"""
Convert PPTX input into a Slidev markdown deck.

Produces full-bleed background slides — each PPTX slide becomes a single
background image. Presenter notes are emitted as clearly-marked TODO
placeholders and MUST be authored manually by reading each rendered slide
image (see the post-conversion checklist in
.github/instructions/slidev.instructions.md).

Usage:
  python scripts/convert-pptx.py <workshop-folder-name>
  python scripts/convert-pptx.py --parts <workshop-folder-name>

  Example:
    python scripts/convert-pptx.py copilot-dev-foundations
    python scripts/convert-pptx.py --parts copilot-dev-foundations

  The script expects either a single PPTX file at:
    source/pptx/<workshop-folder-name>.pptx
  Or split PPTX files at:
    source/pptx/<workshop-folder-name>-part-1.pptx
    source/pptx/<workshop-folder-name>-part-2.pptx

  It produces:
    workshops/<workshop>/<workshop>.slidev.md  — Slidev deck (overwrites)
    public/images/<workshop>/slide-NN.png      — Extracted slide images
"""

import argparse
import sys
import io
import re
import hashlib
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

try:
    from PIL import Image
except ImportError:
    Image = None

# --- Configuration ---
ROOT = Path(__file__).resolve().parent.parent
WORKSHOPS_DIR = ROOT / "workshops"
PUBLIC_DIR = ROOT / "public" / "images"
SOURCE_DIR = ROOT / "source" / "pptx"
WATERMARK_TEXT = "".join(chr(code) for code in (78, 111, 116, 101, 98, 111, 111, 107, 76, 77))


def remove_footers(prs):
    """Remove headers, footers, and generator watermarks from all slides."""
    try:
        # Get slide dimensions from presentation
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        for slide in prs.slides:
            shapes_to_remove = []
            
            for shape in slide.shapes:
                # Remove footer placeholders
                if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                    phf = shape.placeholder_format
                    if phf.type in (14, 15, 16):  # Footer, slide number, date placeholders
                        shapes_to_remove.append(shape)
                
                if hasattr(shape, "text") and WATERMARK_TEXT in shape.text:
                    shapes_to_remove.append(shape)
                
                # Remove small shapes at bottom-right that might be the logo
                # (positioned near the generator watermark text)
                if hasattr(shape, "left") and hasattr(shape, "top"):
                    # Check if shape is in bottom-right area (roughly last 30% x 15%)
                    if (shape.left > slide_width * 0.7 and 
                        shape.top > slide_height * 0.85 and
                        shape.width < slide_width * 0.25):
                        # Likely a small logo/image at bottom-right
                        shapes_to_remove.append(shape)
            
            # Remove identified shapes
            for shape in shapes_to_remove:
                sp = shape.element
                sp.getparent().remove(sp)
    except Exception as e:
        print(f"  Warning: Could not remove footers/watermarks: {e}")


def extract_largest_image(slide):
    """Extract the largest image from a slide (by byte size)."""
    best = None
    best_size = 0

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            blob = shape.image.blob
            if len(blob) > best_size:
                best = {'blob': blob, 'ext': shape.image.ext}
                best_size = len(blob)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    blob = child.image.blob
                    if len(blob) > best_size:
                        best = {'blob': blob, 'ext': child.image.ext}
                        best_size = len(blob)

    return best


def extract_notes(slide):
    """Extract speaker notes text from a slide (if any)."""
    if slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame and notes_frame.text.strip():
            return notes_frame.text.strip()
    return ""


def remove_watermark(img):
    """Remove the generator watermark from the bottom-right corner of an image.

    Copies a horizontal strip of pixels from immediately left of the watermark
    region and pastes it over the watermark. This preserves any underlying
    background color or texture (solid, gradient, graph paper, etc.) and avoids
    single-color fill artifacts on non-uniform backgrounds.

    Only processes 1376x768 images (the standard export size).
    """
    # Only 1376x768 slides carry the generator watermark
    if img.size != (1376, 768):
        return img

    # Watermark bounding box (pixel-exact, with generous margin)
    WM_LEFT, WM_TOP, WM_RIGHT, WM_BOTTOM = 1224, 735, 1376, 766
    WM_W = WM_RIGHT - WM_LEFT

    # Source patch: same y range, directly left of the watermark
    SRC_LEFT = WM_LEFT - WM_W
    SRC_RIGHT = WM_LEFT

    img = img.copy()
    patch = img.crop((SRC_LEFT, WM_TOP, SRC_RIGHT, WM_BOTTOM))
    img.paste(patch, (WM_LEFT, WM_TOP))

    return img


def save_image(blob, ext, slide_index, images_dir):
    """Save an image blob and return the filename."""
    content_hash = hashlib.md5(blob).hexdigest()[:8]
    filename = f"slide-{slide_index:02d}-{content_hash}.{ext}"
    filepath = images_dir / filename

    if Image and ext in ('png', 'jpg', 'jpeg'):
        try:
            img = Image.open(io.BytesIO(blob))
            if img.width > 1920:
                ratio = 1920 / img.width
                new_size = (1920, int(img.height * ratio))
                img = img.resize(new_size, Image.LANCZOS)
            img = remove_watermark(img)
            img.save(filepath, optimize=True)
            return filename
        except Exception:
            pass

    filepath.write_bytes(blob)
    return filename

def parse_workshop_sections(workshop_path):
    """Parse the workshop markdown file into numbered sections with content.

    Returns a list of dicts: [{'title': str, 'content': str}, ...]
    """
    if not workshop_path.exists():
        return []

    text = workshop_path.read_text(encoding='utf-8')

    # Split on H2 headers that start with a number: ## 1. Title (XX min)
    section_pattern = re.compile(r'^## (\d+)\. (.+?)$', re.MULTILINE)
    matches = list(section_pattern.finditer(text))

    if not matches:
        return []

    sections = []
    for i, match in enumerate(matches):
        title = match.group(2).strip()
        start = match.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        content = text[start:end].strip()
        sections.append({'title': title, 'content': content})

    return sections


def extract_talking_points(section_content):
    """Extract key talking points from a workshop section for presenter notes.

    Pulls from Key Points, Safety Moments, and Discussion Points subsections.
    Returns a condensed talk-track string.
    """
    points = []

    # Extract Key Points bullets
    key_points_match = re.search(
        r'### Key Points\s*\n(.*?)(?=\n###|\n---|\Z)',
        section_content, re.DOTALL
    )
    if key_points_match:
        bullets = re.findall(r'^- (.+?)$', key_points_match.group(1), re.MULTILINE)
        # Take the most important bullets (first 3-4)
        for b in bullets[:4]:
            # Strip markdown formatting
            clean = re.sub(r'\*\*(.+?)\*\*', r'\1', b)
            clean = re.sub(r'`(.+?)`', r'\1', clean)
            points.append(clean.strip())

    # Extract Safety Moment
    safety_match = re.search(
        r'### .+?Safety Moment\s*\n(.*?)(?=\n###|\n---|\Z)',
        section_content, re.DOTALL
    )
    if safety_match:
        bullets = re.findall(r'^- (.+?)$', safety_match.group(1), re.MULTILINE)
        if bullets:
            clean = re.sub(r'\*\*(.+?)\*\*', r'\1', bullets[0])
            points.append(f"Safety: {clean.strip()}")

    # Extract Discussion Points (just first one as audience hook)
    discussion_match = re.search(
        r'### Discussion Points\s*\n(.*?)(?=\n###|\n---|\Z)',
        section_content, re.DOTALL
    )
    if discussion_match:
        bullets = re.findall(r'^- (.+?)$', discussion_match.group(1), re.MULTILINE)
        if bullets:
            points.append(f"Ask the audience: {bullets[0].strip()}")

    return points


def generate_notes_from_workshop(workshop_path, num_slides):
    """Generate presenter notes for each slide from the workshop source file.

    Distributes workshop sections across slides proportionally. The first
    slide gets a cover note, the last slide gets a wrap-up note.
    """
    sections = parse_workshop_sections(workshop_path)
    if not sections:
        return ["" for _ in range(num_slides)]

    notes = [""] * num_slides

    # Generate cover slide note from workshop overview
    text = workshop_path.read_text(encoding='utf-8')
    overview_match = re.search(
        r'## Workshop Overview\s*\n(.*?)(?=\n##|\n---)',
        text, re.DOTALL
    )
    if overview_match:
        overview = overview_match.group(1).strip()
        # Take first paragraph as cover note
        first_para = overview.split('\n\n')[0].strip()
        # Remove markdown formatting
        first_para = re.sub(r'\*\*(.+?)\*\*', r'\1', first_para)
        notes[0] = f"Welcome and frame the module. {first_para}"

    # Remaining slides distributed across sections
    content_slides = num_slides - 1  # exclude cover
    slides_per_section = max(1, content_slides / len(sections))

    for sec_idx, section in enumerate(sections):
        # Which slide range does this section map to?
        start_slide = 1 + int(sec_idx * slides_per_section)
        end_slide = 1 + int((sec_idx + 1) * slides_per_section)
        end_slide = min(end_slide, num_slides)

        points = extract_talking_points(section['content'])
        if not points:
            continue

        # First slide of the section gets the bulk of the notes
        section_intro = f"Section: {section['title']}. "
        first_slide_points = points[:3]
        note_text = section_intro + " ".join(first_slide_points)
        if start_slide < num_slides:
            notes[start_slide] = note_text

        # Additional slides in this section get remaining points
        remaining_points = points[3:]
        for offset, point in enumerate(remaining_points):
            slide_idx = start_slide + 1 + offset
            if slide_idx < end_slide and slide_idx < num_slides:
                notes[slide_idx] = point

    return notes


def generate_slidev(image_files, notes_list, workshop_name, title):
    """Generate a Slidev markdown deck with full-bleed image slides.

    Uses background: frontmatter so that Slidev's resolveAssetUrl() prepends
    the --base path at runtime, which is required for GitHub Pages subpath
    deployments. The image-full layout renders the background image full-bleed
    via handleBackground().
    """
    lines = [
        "---",
        "theme: ../../themes/github",
        f'title: "{title}"',
        "info: |",
        f"  Generated from PPTX presentation for {workshop_name}",
        f'ghFooterTitle: "{title}"',
        'ghFooterLabel: ""',
        "drawings:",
        "  persist: false",
        "transition: slide-left",
        "mdc: true",
        "layout: image-full",
        f"background: /images/{workshop_name}/{image_files[0]}",
        "---",
        "",
    ]

    # Cover slide notes (index 0)
    cover_note = notes_list[0] if notes_list and notes_list[0] else "Presenter notes for cover slide"
    if not cover_note.startswith("<!--"):
        cover_note = f"<!-- {cover_note} -->"
    lines.append(cover_note)
    lines.append("")

    for i, img_file in enumerate(image_files[1:], start=1):
        notes = notes_list[i] if i < len(notes_list) and notes_list[i] else "<!-- Presenter notes -->"
        if not notes.startswith("<!--"):
            notes = f"<!-- {notes} -->"

        lines.append("---")
        lines.append("layout: image-full")
        lines.append(f"background: /images/{workshop_name}/{img_file}")
        lines.append("---")
        lines.append("")
        lines.append(notes)
        lines.append("")

    return "\n".join(lines)


def part_sort_key(path):
    """Sort part files by trailing part number, then filename."""
    match = re.search(r'-part-(\d+)\.pptx$', path.name, re.IGNORECASE)
    if match:
        return (int(match.group(1)), path.name.lower())
    return (9999, path.name.lower())


def resolve_pptx_paths(workshop_name, explicit_files, parts_only):
    """Resolve single or split PPTX inputs for a workshop."""
    if explicit_files:
        paths = []
        for value in explicit_files:
            path = Path(value)
            if not path.is_absolute():
                path = SOURCE_DIR / value
            paths.append(path)
        return paths

    if parts_only:
        return sorted(SOURCE_DIR.glob(f"{workshop_name}-part-*.pptx"), key=part_sort_key)

    single_path = SOURCE_DIR / f"{workshop_name}.pptx"
    if single_path.exists():
        return [single_path]

    return sorted(SOURCE_DIR.glob(f"{workshop_name}-part-*.pptx"), key=part_sort_key)


def main():
    parser = argparse.ArgumentParser(
        description="Convert single or split PPTX inputs into one Slidev deck."
    )
    parser.add_argument(
        "--parts",
        action="store_true",
        help="Use source/pptx/<workshop-folder-name>-part-*.pptx inputs.",
    )
    parser.add_argument("workshop_name")
    parser.add_argument(
        "pptx_files",
        nargs="*",
        help="Optional PPTX filenames or paths to combine in the provided order.",
    )
    args = parser.parse_args()

    workshop_name = args.workshop_name
    pptx_paths = resolve_pptx_paths(workshop_name, args.pptx_files, args.parts)

    if not pptx_paths:
        print(f"ERROR: No PPTX files found for: {workshop_name}")
        print(f"  Single input: source/pptx/{workshop_name}.pptx")
        print(f"  Split inputs: source/pptx/{workshop_name}-part-1.pptx, source/pptx/{workshop_name}-part-2.pptx")
        sys.exit(1)

    missing_paths = [path for path in pptx_paths if not path.exists()]
    if missing_paths:
        print("ERROR: PPTX file not found:")
        for path in missing_paths:
            print(f"  {path}")
        sys.exit(1)

    print(f"Workshop:   {workshop_name}")
    print("Inputs:")
    for path in pptx_paths:
        print(f"  - {path.relative_to(ROOT) if path.is_relative_to(ROOT) else path}")

    # Ensure output directories exist
    workshop_dir = WORKSHOPS_DIR / workshop_name
    images_dir = PUBLIC_DIR / workshop_name
    workshop_dir.mkdir(parents=True, exist_ok=True)
    images_dir.mkdir(parents=True, exist_ok=True)

    # Remove old slide images from this workshop
    for old_file in images_dir.glob("slide-*"):
        old_file.unlink()

    image_files = []
    notes_list = []
    title = workshop_name.replace("-", " ").title()
    total_slides = 0
    slide_index = 1

    for deck_index, pptx_path in enumerate(pptx_paths, start=1):
        prs = Presentation(str(pptx_path))
        total_slides += len(prs.slides)
        print(f"Deck {deck_index}:   {pptx_path.name} ({len(prs.slides)} slides)")

        remove_footers(prs)

        if deck_index == 1:
            for slide in prs.slides:
                if slide.shapes.title and slide.shapes.title.text.strip():
                    title = slide.shapes.title.text.strip()
                    break

        for slide in prs.slides:
            img_data = extract_largest_image(slide)
            notes = extract_notes(slide)
            notes_list.append(notes)

            if img_data:
                filename = save_image(img_data['blob'], img_data['ext'], slide_index, images_dir)
                image_files.append(filename)
                print(f"  Slide {slide_index:2d}: {filename}" + (f" [notes]" if notes else ""))
            else:
                print(f"  Slide {slide_index:2d}: WARNING — no image found, skipping")
            slide_index += 1

    if not image_files:
        print("ERROR: No images extracted from PPTX")
        sys.exit(1)

    print(f"Slides:     {total_slides}")
    print(f"Title:      {title}")

    # Generate placeholder notes for slides with no authored PPTX speaker notes.
    # We deliberately do NOT infer notes from the workshop file by slide index —
    # that produced misaligned, wrong notes. Notes are authored manually against
    # each rendered slide image (see slidev.instructions.md post-conversion checklist).
    if not has_pptx_notes:
        print("Notes:      Emitting TODO placeholders — author notes per slide image")
    notes_list = [n if n else NOTES_TODO_PLACEHOLDER for n in notes_list]

    # Generate Slidev markdown
    md_content = generate_slidev(image_files, notes_list, workshop_name, title)

    # Write output
    output_path = workshop_dir / f"{workshop_name}.slidev.md"
    output_path.write_text(md_content, encoding='utf-8')
    print(f"\nOutput:     {output_path}")
    print(f"\nConversion complete!")
    print(f"  - {len(image_files)} slides with background images")
    print(f"  - Replace every TODO placeholder: author notes from each slide image")
    print(f"    (see slidev.instructions.md post-conversion checklist)")
    print(f"  - Preview: npx slidev {output_path.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
